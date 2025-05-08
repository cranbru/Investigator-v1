import os
import sys
import threading
import webbrowser
import string
import ctypes
import base64
import mimetypes
import time
import platform
import io
from flask import Flask, render_template, request, jsonify, send_file
from datetime import datetime

try:
    import docx
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False
    print("python-docx not available, Word files will not be previewed")

try:
    import pptx
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    print("python-pptx not available, PowerPoint files will not be previewed")

class Investigator:
    def __init__(self, allowed_dir=None):
        self.allowed_dir = os.path.abspath(allowed_dir) if allowed_dir else None
        self.log = []
        self.recent_operations = {}  
        self.last_media_access = {}  
       
        mimetypes.init()
        
        self.target_files = self.load_target_files()

    def load_target_files(self):
        target_list = set()
        try:
            script_dir = os.path.dirname(os.path.abspath(__file__))
            target_file_path = os.path.join(script_dir, 'target.txt')
            print(f"Looking for target file at: {target_file_path}")
            
            if os.path.exists(target_file_path):
                with open(target_file_path, 'r', encoding='utf-8') as f:
                    for line in f:
                        line = line.strip()
                        if line and not line.startswith('#'):
                            target_list.add(line.lower())
                print(f"Loaded {len(target_list)} target filenames for highlighting")
                
                sample_entries = list(target_list)[:5] if len(target_list) >= 5 else list(target_list)
                print(f"Sample entries: {sample_entries}")
            else:
                print(f"Target file not found at: {target_file_path}")
        except Exception as e:
            print(f"Error loading target file: {e}")
            import traceback
            traceback.print_exc()
        
        return target_list

    def is_target_file(self, filename):
        if not filename or not self.target_files:
            return False
            
        filename_lower = filename.lower()
        
        if filename_lower in self.target_files:
            print(f"Target match (exact): {filename}")
            return True
        
        if "password.txt" in self.target_files and filename_lower == "passwords.txt":
            print(f"Target match (variation): {filename}")
            return True
        if "passwords.txt" in self.target_files and filename_lower == "password.txt":
            print(f"Target match (variation): {filename}")
            return True
            
        for pattern in self.target_files:
            if not pattern:
                continue
                
            if '*' in pattern:
                if pattern.startswith('*') and pattern.count('*') == 1:
                    suffix = pattern[1:].lower()
                    if suffix and filename_lower.endswith(suffix):
                        print(f"Target match (suffix): {filename} matches {pattern}")
                        return True
                
                elif pattern.endswith('*') and pattern.count('*') == 1:
                    prefix = pattern[:-1].lower()
                    if prefix and filename_lower.startswith(prefix):
                        print(f"Target match (prefix): {filename} matches {pattern}")
                        return True
            
            elif len(pattern) > 3:
                if pattern in filename_lower:
                    print(f"Target match (substring): {filename} contains {pattern}")
                    return True
        
        return False

    def is_within_allowed_dir(self, path):
        if self.allowed_dir is None:
            return True
        return os.path.abspath(path).startswith(self.allowed_dir)

    def add_log_entry(self, action, path, status, message=None, is_media=False):
        timestamp = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        
        clean_path = path.split('?')[0] if '?' in path else path
        
        if is_media:
            current_time = time.time()
            if clean_path in self.last_media_access:
                if current_time - self.last_media_access[clean_path] < 5:
                    return None
            
            self.last_media_access[clean_path] = current_time
            
            self.last_media_access = {k: v for k, v in self.last_media_access.items() 
                                    if current_time - v < 30}
        else:
            operation_key = f"{action}_{clean_path}_{status}"
            
            current_time = time.time()
            if operation_key in self.recent_operations:
                if current_time - self.recent_operations[operation_key] < 2:
                    return None
            
            self.recent_operations[operation_key] = current_time
            
            self.recent_operations = {k: v for k, v in self.recent_operations.items() 
                                    if current_time - v < 5}
        
        log_entry = {
            "timestamp": timestamp,
            "action": action,
            "path": clean_path,
            "status": status
        }
        
        if message:
            log_entry["message"] = message
            
        self.log.append(log_entry)
        return log_entry

    def read_file(self, filepath):
        try:
            mime_type, encoding = mimetypes.guess_type(filepath)
            file_size = os.path.getsize(filepath)
            
            if filepath.lower().endswith('.docx') and DOCX_AVAILABLE:
                content = self.extract_text_from_docx(filepath)
                message = f"[READ] {filepath} (Word document)"
                self.add_log_entry("READ", filepath, "SUCCESS")
                return {
                    "content": content,
                    "type": "text",
                    "is_office": True,
                    "office_type": "word",
                    "mimetype": "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
                }, message
                
            elif filepath.lower().endswith('.pptx') and PPTX_AVAILABLE:
                content = self.extract_text_from_pptx(filepath)
                message = f"[READ] {filepath} (PowerPoint presentation)"
                self.add_log_entry("READ", filepath, "SUCCESS")
                return {
                    "content": content,
                    "type": "text",
                    "is_office": True,
                    "office_type": "powerpoint",
                    "mimetype": "application/vnd.openxmlformats-officedocument.presentationml.presentation"
                }, message
                
            elif self.is_text_file(filepath) and file_size < 5 * 1024 * 1024:
                with open(filepath, 'r', errors='replace') as file:
                    message = f"[READ] {filepath}"
                    self.add_log_entry("READ", filepath, "SUCCESS")
                    return {
                        "content": file.read(),
                        "type": "text",
                        "mimetype": mime_type or "text/plain"
                    }, message
            else:
                message = f"[READ] {filepath}"
                self.add_log_entry("READ", filepath, "SUCCESS")
                return {
                    "content": None,
                    "type": "binary",
                    "mimetype": mime_type or "application/octet-stream",
                    "filename": os.path.basename(filepath),
                    "size": file_size,
                    "last_modified": datetime.fromtimestamp(os.path.getmtime(filepath)).strftime("%Y-%m-%d %H:%M:%S")
                }, message
        except Exception as e:
            message = f"[ERROR] Cannot read: {e}"
            self.add_log_entry("READ", filepath, "ERROR", str(e))
            return {
                "content": None,
                "type": "error",
                "error": str(e)
            }, message

    def get_file_metadata(self, filepath):
        try:
            file_stat = os.stat(filepath)
            filename = os.path.basename(filepath)
            file_size = file_stat.st_size
            mime_type, encoding = mimetypes.guess_type(filepath)
            
            metadata = {
                "basic": {
                    "filename": filename,
                    "filepath": filepath,
                    "size": file_size,
                    "size_human": self.format_size(file_size),
                    "mimetype": mime_type or "unknown",
                    "encoding": encoding,
                    "extension": os.path.splitext(filepath)[1],
                },
                "timestamps": {
                    "created": datetime.fromtimestamp(file_stat.st_ctime).strftime("%Y-%m-%d %H:%M:%S"),
                    "modified": datetime.fromtimestamp(file_stat.st_mtime).strftime("%Y-%m-%d %H:%M:%S"),
                    "accessed": datetime.fromtimestamp(file_stat.st_atime).strftime("%Y-%m-%d %H:%M:%S"),
                },
                "permissions": {
                    "owner_read": bool(file_stat.st_mode & 0o400),
                    "owner_write": bool(file_stat.st_mode & 0o200),
                    "owner_execute": bool(file_stat.st_mode & 0o100),
                    "group_read": bool(file_stat.st_mode & 0o040),
                    "group_write": bool(file_stat.st_mode & 0o020),
                    "group_execute": bool(file_stat.st_mode & 0o010),
                    "other_read": bool(file_stat.st_mode & 0o004),
                    "other_write": bool(file_stat.st_mode & 0o002),
                    "other_execute": bool(file_stat.st_mode & 0o001),
                    "mode_octal": oct(file_stat.st_mode)[-3:],
                },
                "system": {
                    "inode": file_stat.st_ino,
                    "device": file_stat.st_dev,
                    "platform": platform.system(),
                    "platform_version": platform.version(),
                    "python_version": platform.python_version(),
                }
            }
            
            try:
                if mime_type and mime_type.startswith('image/'):
                    from PIL import Image
                    from PIL.ExifTags import TAGS
                    
                    img = Image.open(filepath)
                    exif_data = {}
                    
                    if hasattr(img, '_getexif') and img._getexif():
                        for tag, value in img._getexif().items():
                            if tag in TAGS:
                                exif_data[TAGS[tag]] = str(value)
                    
                    if exif_data:
                        metadata["exif"] = exif_data
            except ImportError:
                metadata["exif_support"] = "PIL library not installed. Install with 'pip install pillow' for EXIF data."
            except Exception as e:
                metadata["exif_error"] = str(e)
            
            self.add_log_entry("METADATA", filepath, "SUCCESS")
            
            return metadata, f"[SUCCESS] Got metadata for {filepath}"
            
        except Exception as e:
            self.add_log_entry("METADATA", filepath, "ERROR", str(e))
            return {"error": str(e)}, f"[ERROR] Failed to get metadata: {e}"
    
    def format_size(self, size_bytes):
        if size_bytes == 0:
            return "0 B"
        
        size_names = ("B", "KB", "MB", "GB", "TB", "PB", "EB", "ZB", "YB")
        i = 0
        while size_bytes >= 1024 and i < len(size_names) - 1:
            size_bytes /= 1024
            i += 1
            
        return f"{size_bytes:.2f} {size_names[i]}"

    def is_text_file(self, filepath):
        mime_type, encoding = mimetypes.guess_type(filepath)
        
        if mime_type:
            return mime_type.startswith('text/') or mime_type in [
                'application/json', 
                'application/xml',
                'application/javascript',
                'application/x-javascript'
            ]
        
        text_extensions = [
            '.txt', '.md', '.csv', '.log', '.py', '.js', '.html', '.css', 
            '.xml', '.json', '.c', '.cpp', '.h', '.java', '.sh', '.bat',
            '.ps1', '.ini', '.cfg', '.conf', '.yaml', '.yml'
        ]
        
        return any(filepath.lower().endswith(ext) for ext in text_extensions)

    def write_file(self, filepath, content):
        message = "[BLOCKED] Write operations are not allowed."
        self.add_log_entry("WRITE", filepath, "BLOCKED")
        return False, message

    def delete_file(self, filepath):
        message = "[BLOCKED] Delete operations are not allowed."
        self.add_log_entry("DELETE", filepath, "BLOCKED")
        return False, message
    
    def get_logs(self):
        return self.log
    
    def reload_target_files(self):
        old_count = len(self.target_files)
        self.target_files = self.load_target_files()
        new_count = len(self.target_files)
        self.add_log_entry("CONFIG", "target.txt", "SUCCESS", f"Reloaded target files: {old_count} → {new_count}")
        return True, f"Reloaded target files: {old_count} → {new_count}"
    
    def list_directory(self, directory=None):
        if directory == "":
            return self.list_drives()
            
        if directory is None:
            if self.allowed_dir:
                directory = self.allowed_dir
            else:
                return self.list_drives()
        
        try:
            items = []
            for item in os.listdir(directory):
                item_path = os.path.join(directory, item)
                item_type = "directory" if os.path.isdir(item_path) else "file"
                
                icon = self.get_file_icon(item_path) if item_type == "file" else None
                
                size = None
                if item_type == "file":
                    try:
                        size = os.path.getsize(item_path)
                    except:
                        pass
                
                is_target = False
                if item_type == "file":
                    is_target = self.is_target_file(item)
                    if is_target:
                        print(f"Flagged target file in directory listing: {item}")
                
                item_data = {
                    "name": item, 
                    "type": item_type, 
                    "path": item_path,
                    "icon": icon,
                    "size": size,
                    "is_target": is_target
                }
                
                items.append(item_data)
                
            print(f"Listed {len(items)} items, {sum(1 for i in items if i.get('is_target'))} are target files")
            return items, f"[SUCCESS] Listed directory: {directory}"
        except Exception as e:
            print(f"Error listing directory: {e}")
            import traceback
            traceback.print_exc()
            return [], f"[ERROR] Cannot list directory: {e}"
    
    def get_file_icon(self, filepath):
        mime_type, encoding = mimetypes.guess_type(filepath)
        if not mime_type:
            return "📄"
            
        if mime_type.startswith('image/'):
            return "🖼️"
        elif mime_type.startswith('video/'):
            return "🎬"
        elif mime_type.startswith('audio/'):
            return "🎵"
        elif mime_type == 'application/pdf':
            return "📕"
        elif mime_type.startswith('text/'):
            return "📝"
        elif 'spreadsheet' in mime_type or mime_type.endswith('excel'):
            return "📊"
        elif 'presentation' in mime_type or mime_type.endswith('powerpoint') or filepath.endswith('.pptx'):
            return "📊"
        elif 'word' in mime_type or filepath.endswith('.docx'):
            return "📄"
        elif 'zip' in mime_type or 'compressed' in mime_type or 'archive' in mime_type:
            return "📦"
        else:
            return "📄"
    
    def list_drives(self):
        drives = []
        if sys.platform == 'win32':
            bitmask = ctypes.windll.kernel32.GetLogicalDrives()
            for letter in string.ascii_uppercase:
                if bitmask & 1:
                    drive_path = f"{letter}:\\"
                    try:
                        drive_name = f"Drive {letter}"
                        try:
                            volumeNameBuffer = ctypes.create_unicode_buffer(1024)
                            fileSystemNameBuffer = ctypes.create_unicode_buffer(1024)
                            ctypes.windll.kernel32.GetVolumeInformationW(
                                ctypes.c_wchar_p(drive_path),
                                volumeNameBuffer,
                                ctypes.sizeof(volumeNameBuffer),
                                None, None, None,
                                fileSystemNameBuffer,
                                ctypes.sizeof(fileSystemNameBuffer)
                            )
                            if volumeNameBuffer.value:
                                drive_name = f"{letter}: {volumeNameBuffer.value}"
                        except:
                            pass
                            
                        drives.append({"name": drive_name, "type": "drive", "path": drive_path})
                    except:
                        pass
                bitmask >>= 1
        else:
            with open('/proc/mounts', 'r') as f:
                for line in f:
                    parts = line.split()
                    if len(parts) > 1:
                        path = parts[1]
                        if path.startswith('/media/') or path.startswith('/mnt/'):
                            name = os.path.basename(path)
                            drives.append({"name": f"Drive {name}", "type": "drive", "path": path})
        
        cwd = os.getcwd()
        drives.append({"name": "Current Working Directory", "type": "directory", "path": cwd})
        
        return drives, "[SUCCESS] Listed available drives"
    
    def set_allowed_dir(self, directory):
        self.allowed_dir = os.path.abspath(directory) if directory else None
        message = f"[SUCCESS] Set allowed directory to: {self.allowed_dir}"
        self.add_log_entry("CONFIG", self.allowed_dir, "SUCCESS")
        return True, message

    def extract_text_from_docx(self, filepath):
        if not DOCX_AVAILABLE:
            return "Word document preview not available (python-docx not installed)"
            
        try:
            doc = Document(filepath)
            full_text = []
            
            for para in doc.paragraphs:
                if para.text:
                    full_text.append(para.text)
            
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text:
                            row_text.append(cell.text)
                    if row_text:
                        full_text.append(" | ".join(row_text))
            
            return "\n\n".join(full_text)
        except Exception as e:
            return f"Error extracting text from Word document: {str(e)}"
    
    def extract_text_from_pptx(self, filepath):
        if not PPTX_AVAILABLE:
            return "PowerPoint preview not available (python-pptx not installed)"
            
        try:
            prs = Presentation(filepath)
            full_text = []
            
            full_text.append("=== PRESENTATION OVERVIEW ===\n")
            
            for i, slide in enumerate(prs.slides):
                slide_text = [f"\n=== SLIDE {i+1} ==="]
                
                if slide.shapes.title:
                    slide_text.append(f"Title: {slide.shapes.title.text}")
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        if shape != slide.shapes.title:
                            slide_text.append(shape.text)
                
                full_text.append("\n".join(slide_text))
            
            return "\n\n".join(full_text)
        except Exception as e:
            return f"Error extracting text from PowerPoint presentation: {str(e)}"

app = Flask(__name__)
blocker = None

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/api/list', methods=['POST'])
def list_dir():
    data = request.json
    directory = data.get('directory', None)
    items, message = blocker.list_directory(directory)
    
    target_files = [item["name"] for item in items if item.get("is_target")]
    if target_files:
        print(f"API response includes target files: {', '.join(target_files)}")
    
    return jsonify({'items': items, 'message': message})

@app.route('/api/read', methods=['POST'])
def read_file():
    data = request.json
    filepath = data.get('filepath')
    content_data, message = blocker.read_file(filepath)
    return jsonify({'data': content_data, 'message': message})

@app.route('/api/metadata', methods=['POST'])
def get_file_metadata():
    data = request.json
    filepath = data.get('filepath')
    metadata, message = blocker.get_file_metadata(filepath)
    return jsonify({'metadata': metadata, 'message': message})

@app.route('/api/download/<path:filepath>')
def download_file(filepath):
    filepath = os.path.normpath(filepath)
    
    blocker.add_log_entry("DOWNLOAD", filepath, "SUCCESS")
    
    return send_file(filepath, as_attachment=True)

@app.route('/api/media/<path:filepath>')
def serve_media(filepath):
    filepath = os.path.normpath(filepath)
    
    mime_type, _ = mimetypes.guess_type(filepath)
    is_media = mime_type and (
        mime_type.startswith('image/') or 
        mime_type.startswith('video/') or
        mime_type.startswith('audio/') or
        mime_type == 'application/pdf'
    )
    
    blocker.add_log_entry("READ", filepath, "SUCCESS", is_media=is_media)
    
    return send_file(filepath)

@app.route('/api/write', methods=['POST'])
def write_file():
    data = request.json
    filepath = data.get('filepath')
    content = data.get('content')
    success, message = blocker.write_file(filepath, content)
    return jsonify({'success': success, 'message': message})

@app.route('/api/delete', methods=['POST'])
def delete_file():
    data = request.json
    filepath = data.get('filepath')
    success, message = blocker.delete_file(filepath)
    return jsonify({'success': success, 'message': message})

@app.route('/api/logs', methods=['GET'])
def get_logs():
    return jsonify({'logs': blocker.get_logs()})

@app.route('/api/set_allowed_dir', methods=['POST'])
def set_allowed_dir():
    data = request.json
    directory = data.get('directory')
    success, message = blocker.set_allowed_dir(directory)
    return jsonify({'success': success, 'message': message})

@app.route('/api/debug/target_files')
def debug_target_files():
    return jsonify({
        'target_files_count': len(blocker.target_files),
        'target_files': list(blocker.target_files)[:20],
        'test_matches': {
            'passwords.txt': blocker.is_target_file('passwords.txt'),
            'todo.txt': blocker.is_target_file('todo.txt'),
            'Todo.txt': blocker.is_target_file('Todo.txt'),
            'resume.docx': blocker.is_target_file('resume.docx'),
            'sample.DAT': blocker.is_target_file('sample.DAT'),
            'normal.txt': blocker.is_target_file('normal.txt'),
            'password.txt': blocker.is_target_file('password.txt'),
            'osintctf.pdf': blocker.is_target_file('osintctf.pdf')
        }
    })

@app.route('/api/reload_targets', methods=['POST'])
def reload_targets():
    success, message = blocker.reload_target_files()
    return jsonify({'success': success, 'message': message})

def open_browser():
    webbrowser.open('http://localhost:5000')

def start_server(allowed_dir=None):
    app.template_folder = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'templates')
    app.static_folder = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'static')
    
    global blocker
    blocker = Investigator(allowed_dir)
    
    threading.Timer(1.5, open_browser).start()
    app.run(host='0.0.0.0', debug=False)

if __name__ == "__main__":
    allowed_directory = None
    if len(sys.argv) > 1:
        allowed_directory = sys.argv[1]
        
    print(f"Starting Investigator v1{' for directory: ' + allowed_directory if allowed_directory else ''}")
    print("Opening web interface...")
    start_server(allowed_directory) 